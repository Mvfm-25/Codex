import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.dml import MSO_THEME_COLOR

BASE_DIR = r"C:\GDM-05\Codex\PUC\ux\trabalhos\t2"
IMG_DIR  = os.path.join(BASE_DIR, "imagensTelas")

# ── Paleta ──────────────────────────────────────────────────────────────────
C_BG_DARK   = RGBColor(0x16, 0x16, 0x2E)
C_PURPLE    = RGBColor(0x5C, 0x4F, 0xFF)
C_PURPLE_LT = RGBColor(0x9B, 0x93, 0xFF)
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT_BG  = RGBColor(0xF5, 0xF5, 0xFF)
C_CARD_BG   = RGBColor(0xFF, 0xFF, 0xFF)
C_GRAY      = RGBColor(0x77, 0x77, 0x88)
C_DARK_TXT  = RGBColor(0x16, 0x16, 0x2E)
C_TEAL      = RGBColor(0x00, 0x96, 0x88)
C_VIOLET    = RGBColor(0x8E, 0x24, 0xAA)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ── Helpers ──────────────────────────────────────────────────────────────────
def new_slide():
    return prs.slides.add_slide(BLANK)

def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, l, t, w, h, fill=None, line=None, line_w=1.0):
    sh = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if line:
        sh.line.color.rgb = line; sh.line.width = Pt(line_w)
    else:
        sh.line.fill.background()
    return sh

def txt(slide, text, l, t, w, h, size=13, bold=False, italic=False,
        color=C_DARK_TXT, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color
    return tb

def header_bar(slide, title, sub=None):
    box(slide, 0, 0, 13.33, 0.85, fill=C_PURPLE)
    txt(slide, title, 0.35, 0.12, 10.0, 0.65,
        size=22, bold=True, color=C_WHITE)
    if sub:
        txt(slide, sub, 0.35, 0.52, 10.0, 0.35,
            size=11, italic=True, color=C_WHITE)

# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 1 — CAPA
# ════════════════════════════════════════════════════════════════════════════
s1 = new_slide()
set_bg(s1, C_BG_DARK)

box(s1, 0,    0,    0.07, 7.5, fill=C_PURPLE)
box(s1, 0.07, 0,    13.26, 0.06, fill=C_PURPLE)

txt(s1, "Plataforma de Qualificação Profissional",
    1.0, 2.0, 11.33, 1.5,
    size=38, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

box(s1, 4.2, 3.65, 4.93, 0.05, fill=C_PURPLE_LT)

txt(s1, "Protótipo de Baixa Fidelidade  ·  Mapeamento por Personas",
    1.0, 3.8, 11.33, 0.55,
    size=15, italic=True, color=C_PURPLE_LT, align=PP_ALIGN.CENTER)

txt(s1, "Experiência do Usuário  ·  2026",
    1.0, 4.4, 11.33, 0.45,
    size=12, color=C_GRAY, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 2 — 6 TELAS IMPLEMENTADAS
# ════════════════════════════════════════════════════════════════════════════
s2 = new_slide()
set_bg(s2, C_LIGHT_BG)
header_bar(s2, "Visão Geral  —  6 Telas Implementadas")

screens = [
    ("01", "Tela Inicial",       "Captura de tela 2026-05-26 135014.png"),
    ("02", "Busca de Cursos",    "Captura de tela 2026-05-26 135040.png"),
    ("03", "Player do Curso",    "Captura de tela 2026-05-26 135019.png"),
    ("04", "Fórum",              "Captura de tela 2026-05-26 135027.png"),
    ("05", "Mentorias",          "Captura de tela 2026-05-26 135033.png"),
    ("06", "Perfil",             "Captura de tela 2026-05-26 135008.png"),
]

for i, (num, title, img_file) in enumerate(screens):
    col = i % 3
    row = i // 3
    l = 0.35 + col * 4.32
    t = 1.05 + row * 3.05
    w, h = 4.0, 2.85

    box(s2, l, t, w, h, fill=C_CARD_BG, line=C_PURPLE_LT, line_w=1.2)
    box(s2, l, t, w, 0.42, fill=C_PURPLE)

    txt(s2, num,   l + 0.12, t + 0.06, 0.4,    0.32, size=14, bold=True, color=C_WHITE)
    txt(s2, title, l + 0.55, t + 0.07, w-0.65, 0.32, size=13, bold=True, color=C_WHITE)

    img_path = os.path.join(IMG_DIR, img_file)
    s2.shapes.add_picture(img_path,
                          Inches(l + 0.08), Inches(t + 0.45),
                          Inches(w - 0.16), Inches(h - 0.53))

# ════════════════════════════════════════════════════════════════════════════
#  SLIDES 3–5 — PERSONAS
# ════════════════════════════════════════════════════════════════════════════
# patterns = lista de nomes de padrões do ui-patterns.com para cada tela
personas = [
    {
        "name":  "Gabriel Costa",
        "role":  "Desenvolvedor Backend · 23 anos",
        "quote": '"Quero dominar o que faço, não só passar por cima"',
        "color": C_PURPLE,
        "items": [
            ("Player do Curso",
             "Sidebar com aulas numeradas permite navegar no próprio ritmo → estudo diário com progressão autônoma.",
             ["Module Tabs", "Progressive Disclosure", "Article List"]),
            ("Tela Inicial",
             "Barra de progresso em 'Em andamento' torna a evolução visível a cada acesso → feedback contínuo.",
             ["Dashboard", "Completeness Meter", "Cards"]),
            ("Fórum",
             "Tópicos com tags técnicas (#Code #Java #Iniciantes) viabilizam troca aprofundada entre pares.",
             ["Tagging", "Activity Stream", "Search Filters"]),
            ("Busca de Cursos",
             "Filtros por área (Front-End, Back-End, Dev-Ops, QA) permitem encontrar conteúdo técnico específico.",
             ["Search Filters", "Categorization", "Cards", "Thumbnail"]),
        ],
    },
    {
        "name":  "Beatriz Mendes",
        "role":  "Estudante em transição · 22 anos",
        "quote": '"Não sei bem pra onde ir, mas sei que preciso mudar"',
        "color": C_VIOLET,
        "items": [
            ("Busca de Cursos",
             "'Cursos Recomendados' + filtros de área oferecem orientação antes do compromisso financeiro.",
             ["Search Filters", "Categorization", "Cards", "Thumbnail"]),
            ("Mentorias",
             "Chat com mentores experientes (badge Mestre) e grupos de estudo fornecem suporte acessível e contínuo.",
             ["Chat", "Activity Stream", "Reputation"]),
            ("Perfil",
             "Histórico de cursos concluídos com datas serve como evidência do esforço e validação do investimento.",
             ["Dashboard", "Archive", "Completeness Meter"]),
            ("Fórum",
             "Comunidade de discussão conecta Beatriz a profissionais da área, reduzindo a insegurança de transição.",
             ["Tagging", "Activity Stream", "Reaction"]),
        ],
    },
    {
        "name":  "Lucas Oliveira",
        "role":  "Candidato competitivo · 21 anos",
        "quote": '"Preciso de resultado concreto, não de mais teoria"',
        "color": C_TEAL,
        "items": [
            ("Player do Curso",
             "Lista de aulas visível e conteúdo sequencial oferecem estrutura clara para atividades práticas contextualizadas.",
             ["Module Tabs", "Progressive Disclosure", "Article List"]),
            ("Mentorias",
             "Acesso imediato a mentor (Mestre) e grupos de estudo = suporte quando travar no conteúdo.",
             ["Chat", "Activity Stream", "Reputation"]),
            ("Perfil",
             "Histórico de conclusões com datas funciona como histórico verificável para currículo e portfólio.",
             ["Dashboard", "Archive", "Completeness Meter"]),
            ("Busca de Cursos",
             "Cursos recomendados com categorias objetivas dão clareza sobre o que será recebido pelo investimento.",
             ["Search Filters", "Categorization", "Cards", "Thumbnail"]),
        ],
    },
]

# Tamanho dos cards ajustado para acomodar a linha de padrões
CARD_H   = 1.78   # altura do card
CARD_GAP = 1.82   # espaçamento entre inícios de cards

C_PILL_BG = RGBColor(0xEE, 0xEB, 0xFF)   # fundo claro para área de padrões (roxo-lavanda)

for persona in personas:
    s = new_slide()
    set_bg(s, C_LIGHT_BG)
    pc = persona["color"]

    # ── Painel esquerdo ──
    box(s, 0, 0, 3.7, 7.5, fill=pc)

    # Círculo avatar (fundo branco)
    av = s.shapes.add_shape(9, Inches(1.15), Inches(0.28), Inches(1.4), Inches(1.4))
    av.fill.solid(); av.fill.fore_color.rgb = C_WHITE
    av.line.fill.background()

    # Ícone de usuário genérico: cabeça + ombros na cor da persona
    head = s.shapes.add_shape(9, Inches(1.64), Inches(0.49), Inches(0.42), Inches(0.42))
    head.fill.solid(); head.fill.fore_color.rgb = pc
    head.line.fill.background()

    body = s.shapes.add_shape(9, Inches(1.50), Inches(1.00), Inches(0.70), Inches(0.55))
    body.fill.solid(); body.fill.fore_color.rgb = pc
    body.line.fill.background()

    txt(s, persona["name"],  0.18, 1.88, 3.35, 0.65,
        size=18, bold=True, color=C_WHITE)
    txt(s, persona["role"],  0.18, 2.55, 3.35, 0.45,
        size=11, italic=True, color=C_WHITE)

    box(s, 0.18, 3.08, 3.34, 0.04, fill=RGBColor(0xFF, 0xFF, 0xFF))

    txt(s, persona["quote"], 0.18, 3.2, 3.35, 1.5,
        size=11, italic=True, color=C_WHITE)

    txt(s, "Objetivos Práticos Atendidos", 0.18, 5.5, 3.35, 0.5,
        size=10, bold=True, color=C_WHITE)

    # ── Cards de objetivos ──
    for j, (tela, desc, patterns) in enumerate(persona["items"]):
        t = 0.18 + j * CARD_GAP

        # Card body
        box(s, 3.9, t, 9.15, CARD_H, fill=C_CARD_BG, line=pc, line_w=1.0)

        # Colored header bar
        box(s, 3.9, t, 9.15, 0.38, fill=pc)
        txt(s, tela, 4.05, t + 0.04, 6.8, 0.32,
            size=12, bold=True, color=C_WHITE)
        txt(s, "✓", 12.65, t + 0.04, 0.3, 0.32,
            size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

        # Description
        txt(s, desc, 4.05, t + 0.44, 8.85, 0.76,
            size=10, color=C_GRAY)

        # ── Padrões UI ──
        pattern_top = t + CARD_H - 0.50
        box(s, 3.9, pattern_top, 9.15, 0.48, fill=C_PILL_BG)
        box(s, 3.9, pattern_top, 9.15, 0.03, fill=pc)  # thin top border

        txt(s, "Padrões UI:", 4.05, pattern_top + 0.06, 1.4, 0.22,
            size=9, bold=True, color=pc)
        pills_text = "  ·  ".join(patterns)
        txt(s, pills_text, 5.35, pattern_top + 0.06, 7.55, 0.38,
            size=9, color=C_DARK_TXT)

# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 6 — MAPEAMENTO TELAS × PERSONAS
# ════════════════════════════════════════════════════════════════════════════
s6 = new_slide()
set_bg(s6, C_LIGHT_BG)
header_bar(s6, "Mapeamento  —  Telas × Personas",
           sub="✓ atende diretamente   ◎ atende parcialmente   — não contemplado")

headers = ["Tela", "Gabriel\n(Técnico)", "Beatriz\n(Transição)", "Lucas\n(Emprego)"]
rows = [
    ("Tela Inicial",       "✓  Progresso diário",    "◎  Recomendações",       "◎  Recomendações"),
    ("Busca de Cursos",    "✓  Filtros técnicos",     "✓  Orientação de área",  "✓  Clareza da oferta"),
    ("Player do Curso",    "✓  Ritmo próprio",        "◎  Consumo de conteúdo", "✓  Prática sequencial"),
    ("Fórum",              "✓  Comunidade técnica",   "✓  Apoio profissional",  "✓  Suporte ao travar"),
    ("Mentorias",          "◎  Troca com pares",      "✓  Mentor direto",       "✓  Mentor disponível"),
    ("Perfil",             "◎  Registro pessoal",     "✓  Validação esforço",   "✓  Currículo/portfólio"),
]

col_w  = [3.0, 3.0, 3.0, 3.0]
col_x  = [0.3, 3.45, 6.6, 9.75]
rh     = 0.88
hdr_t  = 0.98

# Cabeçalho da tabela
for ci in range(4):
    box(s6, col_x[ci], hdr_t, col_w[ci] - 0.06, 0.52, fill=C_BG_DARK)
    txt(s6, headers[ci], col_x[ci]+0.1, hdr_t+0.04, col_w[ci]-0.2, 0.48,
        size=11, bold=True, color=C_WHITE,
        align=PP_ALIGN.LEFT if ci==0 else PP_ALIGN.CENTER)

# Linhas de dados
alt_bg = RGBColor(0xEA, 0xEA, 0xFF)
for ri, row in enumerate(rows):
    t = hdr_t + 0.57 + ri * rh
    fill_c = C_CARD_BG if ri % 2 == 0 else alt_bg
    for ci, (cell, cw, cx) in enumerate(zip(row, col_w, col_x)):
        box(s6, cx, t, cw-0.06, rh-0.06, fill=fill_c, line=C_PURPLE_LT, line_w=0.5)
        if cell.startswith("✓"):
            c = C_PURPLE; bld = True
        elif cell.startswith("◎"):
            c = C_TEAL;   bld = False
        else:
            c = C_GRAY;   bld = False
        sz  = 12 if ci == 0 else 11
        bld = bld or (ci == 0)
        aln = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
        txt(s6, cell, cx+0.12, t+0.16, cw-0.24, rh-0.22,
            size=sz, bold=(ci==0), color=c if ci>0 else C_DARK_TXT, align=aln)

# ════════════════════════════════════════════════════════════════════════════
out = os.path.join(BASE_DIR, "apresentacao-ux-v3.pptx")
prs.save(out)
print(f"Salvo em: {out}")
